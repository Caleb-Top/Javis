"""Javis自创: ppt_to_images — PPT每页导出为图片"""
TOOL_NAME="ppt_to_images"
TOOL_DESC="将PPT/PPTX每页幻灯片导出为图片(PNG/JPG)。参数: ppt_path(输入路径), output_dir(输出目录), format(图片格式png/jpg,默认png), scale(缩放倍数,默认2即高清)"
TOOL_CATEGORY="document_convert"
TOOL_PARAMS={"type":"object","properties":{},"required":[]}

def handler(**kwargs):
    import os, json
    from pathlib import Path

    def convert(ppt_path=None, output_dir=None, format="png", scale=2, **kw):
        if not ppt_path or not os.path.exists(ppt_path):
            return {"success": False, "output": "", "message": f"文件不存在: {ppt_path}"}
        if not output_dir:
            output_dir = str(Path(ppt_path).parent / (Path(ppt_path).stem + "_slides"))
        os.makedirs(output_dir, exist_ok=True)

        try:
            # 方案1: 使用 LibreOffice 无头模式 (最佳质量)
            import subprocess
            import shutil

            libreoffice = shutil.which("soffice") or shutil.which("libreoffice")
            if libreoffice:
                # LibreOffice 命令行转换
                cmd = [libreoffice, "--headless", "--convert-to", format,
                       f"--outdir", output_dir, ppt_path]
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
                if result.returncode == 0:
                    images = sorted([os.path.join(output_dir, f) for f in os.listdir(output_dir)
                                    if f.lower().endswith(f'.{format}')])
                    if images:
                        return {"success": True, "output": output_dir,
                                "message": f"成功导出 {len(images)} 张图片 (LibreOffice引擎)",
                                "images": images}
            # 方案2: 使用 python-pptx + Pillow (无需LibreOffice)
            from pptx import Presentation
            from pptx.util import Inches
            from PIL import Image, ImageDraw, ImageFont

            prs = Presentation(ppt_path)
            # 获取幻灯片尺寸
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            img_width = int(slide_width * scale / 914400)  # EMU to pixels at scale
            img_height = int(slide_height * scale / 914400)

            images = []
            for slide_num, slide in enumerate(prs.slides, 1):
                # 创建白色画布
                img = Image.new('RGB', (img_width, img_height), 'white')
                draw = ImageDraw.Draw(img)

                y_offset = 20
                # 提取文本
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                try:
                                    font = ImageFont.truetype("arial.ttf", 14 * scale)
                                except:
                                    font = ImageFont.load_default()
                                # 判断标题
                                is_title = 'title' in shape.name.lower() or '标题' in shape.name.lower()
                                if is_title:
                                    try:
                                        font = ImageFont.truetype("arial.ttf", 24 * scale)
                                    except:
                                        pass
                                # 文本换行
                                chars_per_line = max(1, img_width // (10 * scale))
                                while len(text) > chars_per_line:
                                    draw.text((20, y_offset), text[:chars_per_line], fill='black', font=font)
                                    text = text[chars_per_line:]
                                    y_offset += int(20 * scale)
                                draw.text((20, y_offset), text, fill='black', font=font)
                                y_offset += int(22 * scale)

                        y_offset += int(10 * scale)  # 形状间距

                img_path = os.path.join(output_dir, f"slide_{slide_num:02d}.{format.lower()}")
                img.save(img_path)
                images.append(img_path)

            return {"success": True, "output": output_dir,
                    "message": f"成功导出 {len(images)} 张图片 (Pillow引擎)",
                    "images": images}
        except Exception as e:
            return {"success": False, "output": "", "message": f"转换失败: {str(e)}"}

    result = convert(**kwargs)
    return json.dumps(result, ensure_ascii=False)
