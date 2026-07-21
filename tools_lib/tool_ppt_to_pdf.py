"""Javis自创: ppt_to_pdf"""
TOOL_NAME="ppt_to_pdf"
TOOL_DESC="将PPT/PPTX演示文稿转换为PDF。参数: ppt_path(输入路径), output_path(可选输出路径)"
TOOL_CATEGORY="document_convert"
TOOL_PARAMS={"type":"object","properties":{},"required":[]}

def handler(**kwargs):
    import os, sys, importlib, json
    from pathlib import Path


    def convert(ppt_path=None, output_path=None, **kw):
        if not ppt_path or not os.path.exists(ppt_path):
            return {"success": False, "output": "", "message": f"文件不存在: {ppt_path}"}
        if not output_path:
            output_path = str(Path(ppt_path).with_suffix('.pdf'))
        try:
            from pptx import Presentation
            from reportlab.lib.pagesizes import A4, landscape
            from reportlab.pdfgen import canvas
            from reportlab.lib.units import mm

            prs = Presentation(ppt_path)
            c = canvas.Canvas(output_path, pagesize=landscape(A4))
            width, height = landscape(A4)

            for slide_num, slide in enumerate(prs.slides, 1):
                y = height - 30*mm
                c.setFont("Helvetica-Bold", 16)
                c.drawString(20*mm, y, f"Slide {slide_num}")
                y -= 15*mm

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text and y > 20*mm:
                                c.setFont("Helvetica", 11)
                                c.drawString(25*mm, y, text[:90])
                                y -= 7*mm

                if slide_num < len(prs.slides):
                    c.showPage()

            c.save()
            return {"success": True, "output": output_path, "message": f"PPT已转换为PDF: {output_path}"}
        except Exception as e:
            return {"success": False, "output": "", "message": f"转换失败: {str(e)}"}
    

    result = convert(**kwargs)
    return json.dumps(result, ensure_ascii=False)
