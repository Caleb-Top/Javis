"""Javis自创: extract_text"""
TOOL_NAME="extract_text"
TOOL_DESC="从任意文档格式(PDF/Word/Excel/PPT/图片/Markdown/TXT)中提取文本内容。万能文本提取器。参数: file_path(文件路径), output_path(可选输出txt路径)"
TOOL_CATEGORY="document_convert"
TOOL_PARAMS={"type":"object","properties":{},"required":[]}

def handler(**kwargs):
    import os, sys, importlib, json
    from pathlib import Path


    def convert(file_path=None, output_path=None, **kw):
        if not file_path or not os.path.exists(file_path):
            return {"success": False, "output": "", "message": f"文件不存在: {file_path}"}

        ext = os.path.splitext(file_path)[1].lower()
        text = ""
        source_type = ""

        try:
            if ext == '.pdf':
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    text = "\n".join([page.extract_text() or "" for page in pdf.pages])
                source_type = "PDF"
            elif ext in ['.docx', '.doc']:
                from docx import Document
                doc = Document(file_path)
                text = "\n".join([p.text for p in doc.paragraphs])
                # 也提取表格
                for table in doc.tables:
                    for row in table.rows:
                        text += "\n" + " | ".join([cell.text for cell in row.cells])
                source_type = "Word"
            elif ext in ['.xlsx', '.xls']:
                import pandas as pd
                dfs = pd.read_excel(file_path, sheet_name=None)
                for sn, df in dfs.items():
                    text += f"\n=== Sheet: {sn} ===\n"
                    text += df.to_string(index=False) + "\n"
                source_type = "Excel"
            elif ext in ['.pptx', '.ppt']:
                from pptx import Presentation
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                text += para.text + "\n"
                source_type = "PPT"
            elif ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
                try:
                    import pytesseract
                    from PIL import Image
                    img = Image.open(file_path)
                    text = pytesseract.image_to_string(img, lang='chi_sim+eng')
                    source_type = "Image(OCR)"
                except:
                    text = "[OCR不可用，请安装pytesseract并配置tesseract路径]"
            elif ext in ['.md', '.markdown']:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
                source_type = "Markdown"
            elif ext == '.txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
                source_type = "Text"
            else:
                # 尝试作为纯文本读取
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        text = f.read()
                    source_type = "Unknown(Binary)"
                except:
                    return {"success": False, "output": "", "message": f"不支持的文件格式: {ext}"}

            if output_path:
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(text)
                return {"success": True, "output": output_path, "message": f"已从{source_type}提取文本", "text_length": len(text)}
            else:
                # 返回前500字符预览
                return {"success": True, "output": "", "message": f"已从{source_type}提取文本", "text_length": len(text), "preview": text[:500]}

        except Exception as e:
            return {"success": False, "output": "", "message": f"提取失败: {str(e)}"}
    

    result = convert(**kwargs)
    return json.dumps(result, ensure_ascii=False)
