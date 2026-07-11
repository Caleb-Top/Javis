"""Javis自创: ppt_to_docx — PPT演示文稿转Word文档"""
TOOL_NAME="ppt_to_docx"
TOOL_DESC="将PPT/PPTX演示文稿转换为Word文档。提取每页标题、文本内容、备注，保留幻灯片结构。参数: ppt_path(输入路径), output_path(可选输出路径), include_notes(是否包含备注,默认true)"
TOOL_CATEGORY="document_convert"
TOOL_PARAMS={"type":"object","properties":{},"required":[]}

def handler(**kwargs):
    import os, json
    from pathlib import Path

    def convert(ppt_path=None, output_path=None, include_notes=True, **kw):
        if not ppt_path or not os.path.exists(ppt_path):
            return {"success": False, "output": "", "message": f"文件不存在: {ppt_path}"}
        if not output_path:
            output_path = str(Path(ppt_path).with_suffix('.docx'))

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from docx import Document
            from docx.shared import Pt as DocxPt, Inches as DocxInches, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            prs = Presentation(ppt_path)
            doc = Document()

            # 标题页
            title = doc.add_heading(f'PPT文档: {Path(ppt_path).stem}', 0)

            for slide_num, slide in enumerate(prs.slides, 1):
                # 每页幻灯片标题
                doc.add_heading(f'幻灯片 {slide_num}', level=1)

                # 提取所有形状中的文本
                shapes_with_text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_content = []
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                text_content.append(t)
                        if text_content:
                            shapes_with_text.append({
                                'name': shape.name,
                                'type': str(shape.shape_type),
                                'text': text_content,
                                'top': shape.top,
                                'left': shape.left
                            })

                # 按位置排序（从上到下）
                shapes_with_text.sort(key=lambda s: (s['top'], s['left']))

                # 写入Word
                for shape_info in shapes_with_text:
                    # 判断是否是标题形状
                    is_title = any(kw in shape_info['name'].lower() for kw in ['title', '标题'])
                    
                    for i, text_line in enumerate(shape_info['text']):
                        if is_title and i == 0:
                            p = doc.add_paragraph()
                            run = p.add_run(text_line)
                            run.bold = True
                            run.font.size = DocxPt(14)
                        else:
                            # 列表项
                            if text_line.startswith('•') or text_line.startswith('-') or text_line.startswith('·'):
                                doc.add_paragraph(text_line, style='List Bullet')
                            elif len(text_line) < 10:
                                p = doc.add_paragraph()
                                run = p.add_run(text_line)
                                run.bold = True
                            else:
                                doc.add_paragraph(text_line)

                # 提取表格
                for shape in slide.shapes:
                    if shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            table_data.append([cell.text for cell in row.cells])
                        if table_data:
                            doc_table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
                            doc_table.style = 'Light Grid Accent 1'
                            for i, row in enumerate(table_data):
                                for j, cell_text in enumerate(row):
                                    doc_table.rows[i].cells[j].text = cell_text
                            doc.add_paragraph()

                # 备注
                if include_notes and slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        p = doc.add_paragraph()
                        run = p.add_run(f'📝 备注: {notes}')
                        run.italic = True
                        run.font.size = DocxPt(9)
                        run.font.color.rgb = RGBColor(128, 128, 128)

                # 分页
                if slide_num < len(prs.slides):
                    doc.add_page_break()

            doc.save(output_path)
            return {"success": True, "output": output_path, "message": f"PPT已转换为Word ({len(prs.slides)}页幻灯片): {output_path}"}
        except Exception as e:
            return {"success": False, "output": "", "message": f"转换失败: {str(e)}"}

    result = convert(**kwargs)
    return json.dumps(result, ensure_ascii=False)
