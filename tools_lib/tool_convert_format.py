"""Javis自创: convert_format — 万能文件格式转换器 v2.0"""
TOOL_NAME="convert_format"
TOOL_DESC="万能文件格式转换器。支持: PDF↔Word↔Excel↔PPT↔图片↔Markdown↔TXT↔HTML↔EPUB↔CSV之间互转, OCR图片文字提取, CAT翻译格式解析(TMX/TBX/XLIFF)。参数: source_path(源文件路径), target_format(目标格式), output_path(可选输出路径)"
TOOL_CATEGORY="document_convert"
TOOL_PARAMS={"type":"object","properties":{},"required":[]}

def handler(**kwargs):
    import os, json
    from pathlib import Path

    def convert(source_path=None, target_format=None, output_path=None, **kw):
        if not source_path or not os.path.exists(source_path):
            return {"success": False, "output": "", "message": f"文件不存在: {source_path}"}
        if not target_format:
            return {"success": False, "output": "", "message": "请指定目标格式: pdf/docx/xlsx/pptx/png/jpg/md/txt/csv/json/html"}

        target_format = target_format.lower().replace('.', '')
        ext = os.path.splitext(source_path)[1].lower().replace('.', '')

        # 规范化格式名
        fmt_alias = {
            'word':'docx','excel':'xlsx','markdown':'md','text':'txt',
            'jpeg':'jpg','powerpoint':'pptx','csv':'csv','html':'html',
            'epub':'epub','json':'json'
        }
        target_format = fmt_alias.get(target_format, target_format)

        if not output_path:
            output_path = str(Path(source_path).parent / (Path(source_path).stem + f".{target_format}"))

        try:
            # ============ PDF → * ============
            if ext == 'pdf':
                if target_format in ['docx']:
                    from pdf2docx import Converter
                    cv = Converter(source_path)
                    cv.convert(output_path)
                    cv.close()
                elif target_format in ['xlsx']:
                    import pdfplumber, pandas as pd
                    with pdfplumber.open(source_path) as pdf:
                        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                            table_count = 0
                            for i, page in enumerate(pdf.pages):
                                tables = page.extract_tables()
                                for j, t in enumerate(tables):
                                    if t and len(t) > 1:
                                        pd.DataFrame(t[1:], columns=t[0]).to_excel(
                                            writer, sheet_name=f"P{i+1}_T{j+1}", index=False)
                                        table_count += 1
                            if table_count == 0:
                                text = "\n".join([p.extract_text() or "" for p in pdf.pages])
                                pd.DataFrame({"Text": text.split('\n')}).to_excel(writer, sheet_name="Text", index=False)
                elif target_format in ['md']:
                    # 委托给专用工具
                    from tool_pdf_to_markdown import handler as md_handler
                    result_str = md_handler(pdf_path=source_path, output_path=output_path)
                    result = json.loads(result_str)
                    return result
                elif target_format in ['png', 'jpg']:
                    import fitz
                    doc = fitz.open(source_path)
                    if doc.page_count == 1:
                        page = doc[0]
                        pix = page.get_pixmap(dpi=200)
                        pix.save(output_path)
                    else:
                        out_dir = str(Path(output_path).parent / Path(output_path).stem)
                        os.makedirs(out_dir, exist_ok=True)
                        for i, page in enumerate(doc):
                            pix = page.get_pixmap(dpi=200)
                            pix.save(os.path.join(out_dir, f"page_{i+1}.{target_format}"))
                        output_path = out_dir
                    doc.close()
                elif target_format in ['txt']:
                    import fitz
                    doc = fitz.open(source_path)
                    text = "\n".join([page.get_text() for page in doc])
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(text)
                    doc.close()
                elif target_format in ['pptx']:
                    import fitz
                    from pptx import Presentation
                    from pptx.util import Inches
                    prs = Presentation()
                    doc = fitz.open(source_path)
                    for page in doc:
                        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
                        pix = page.get_pixmap(dpi=150)
                        img_path = output_path.replace('.pptx', f'_tmp_{page.number}.png')
                        pix.save(img_path)
                        slide.shapes.add_picture(img_path, 0, 0, Inches(10), Inches(7.5))
                    prs.save(output_path)
                    # 清理临时文件
                    import glob
                    for tmp in glob.glob(output_path.replace('.pptx', '_tmp_*.png')):
                        try: os.remove(tmp)
                        except: pass
                    doc.close()
                else:
                    return {"success": False, "message": f"不支持 PDF→{target_format}"}

            # ============ DOCX → * ============
            elif ext in ['docx', 'doc']:
                from docx import Document
                doc = Document(source_path)
                if target_format == 'pdf':
                    from reportlab.lib.pagesizes import A4
                    from reportlab.pdfgen import canvas
                    from reportlab.lib.units import mm
                    from reportlab.pdfbase import pdfmetrics
                    from reportlab.pdfbase.ttfonts import TTFont
                    cjk = "Helvetica"
                    for fp, fn in [("C:/Windows/Fonts/simhei.ttf","SimHei"),("C:/Windows/Fonts/msyh.ttc","MSYaHei")]:
                        if os.path.exists(fp):
                            try: pdfmetrics.registerFont(TTFont(fn, fp)); cjk = fn; break
                            except: pass
                    c = canvas.Canvas(output_path, pagesize=A4)
                    w, h = A4; y = h - 25*mm
                    for para in doc.paragraphs:
                        text = para.text.strip()
                        if not text: continue
                        if y < 20*mm: c.showPage(); y = h - 25*mm
                        c.setFont(cjk, 10)
                        while text:
                            c.drawString(20*mm, y, text[:100]); text = text[100:]; y -= 5*mm
                    for table in doc.tables:
                        if y < 40*mm: c.showPage(); y = h - 25*mm
                        for row in table.rows:
                            row_text = " | ".join(cell.text for cell in row.cells)
                            c.setFont(cjk, 8); c.drawString(20*mm, y, row_text[:120]); y -= 5*mm
                    c.save()
                elif target_format == 'xlsx':
                    import pandas as pd
                    with pd.ExcelWriter(output_path, engine='openpyxl') as w:
                        for i, table in enumerate(doc.tables):
                            data = [[c.text for c in r.cells] for r in table.rows]
                            if data: pd.DataFrame(data[1:], columns=data[0]).to_excel(w, sheet_name=f"T{i+1}", index=False)
                        paras = [p.text for p in doc.paragraphs if p.text.strip()]
                        if paras: pd.DataFrame(paras, columns=["Text"]).to_excel(w, sheet_name="Text", index=False)
                elif target_format in ['md', 'txt']:
                    text = "\n".join([p.text for p in doc.paragraphs])
                    if target_format == 'md':
                        text = "# " + Path(source_path).stem + "\n\n" + text
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(text)
                elif target_format in ['png', 'jpg']:
                    from PIL import Image, ImageDraw, ImageFont
                    text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
                    img = Image.new('RGB', (800, max(800, len(text)//2 + 200)), 'white')
                    draw = ImageDraw.Draw(img)
                    try: font = ImageFont.truetype("C:/Windows/Fonts/simhei.ttf", 14)
                    except: font = ImageFont.load_default()
                    y = 10
                    for line in text.split('\n')[:200]:
                        draw.text((20, y), line[:100], fill='black', font=font); y += 18
                    img.save(output_path)
                elif target_format == 'pptx':
                    from pptx import Presentation
                    prs = Presentation()
                    for para in doc.paragraphs:
                        if para.text.strip():
                            slide = prs.slides.add_slide(prs.slide_layouts[1])
                            slide.shapes.title.text = para.text[:100]
                    prs.save(output_path)
                else:
                    return {"success": False, "message": f"不支持 Word→{target_format}"}

            # ============ XLSX → * ============
            elif ext in ['xlsx', 'xls']:
                import pandas as pd
                dfs = pd.read_excel(source_path, sheet_name=None)
                if target_format == 'pdf':
                    from reportlab.lib.pagesizes import A4
                    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
                    from reportlab.lib import colors
                    doc = SimpleDocTemplate(output_path, pagesize=A4)
                    elements = []
                    for sn, df in dfs.items():
                        data = [df.columns.tolist()] + df.astype(str).values.tolist()
                        t = Table(data)
                        t.setStyle(TableStyle([
                            ('BACKGROUND',(0,0),(-1,0),colors.grey),
                            ('TEXTCOLOR',(0,0),(-1,0),colors.whitesmoke),
                            ('GRID',(0,0),(-1,-1),0.5,colors.black),
                        ]))
                        elements.append(t)
                    doc.build(elements)
                elif target_format == 'docx':
                    from docx import Document
                    doc = Document()
                    doc.add_heading(Path(source_path).stem, 0)
                    for sn, df in dfs.items():
                        doc.add_heading(f'工作表: {sn}', level=1)
                        table = doc.add_table(rows=len(df)+1, cols=len(df.columns))
                        table.style = 'Light Grid Accent 1'
                        for j, col in enumerate(df.columns): table.rows[0].cells[j].text = str(col)
                        for i in range(len(df)):
                            for j in range(len(df.columns)): table.rows[i+1].cells[j].text = str(df.iloc[i,j])
                    doc.save(output_path)
                elif target_format == 'md':
                    # 委托专用工具
                    from tool_xlsx_to_markdown import handler as xm_handler
                    result_str = xm_handler(excel_path=source_path, output_path=output_path)
                    return json.loads(result_str)
                elif target_format == 'txt':
                    with open(output_path, 'w', encoding='utf-8') as f:
                        for sn, df in dfs.items():
                            f.write(f"\n=== {sn} ===\n")
                            f.write(df.to_string(index=False))
                            f.write("\n")
                elif target_format == 'csv':
                    # 取第一个工作表
                    first_sheet = list(dfs.keys())[0]
                    dfs[first_sheet].to_csv(output_path, index=False, encoding='utf-8-sig')
                else:
                    return {"success": False, "message": f"不支持 Excel→{target_format}"}

            # ============ PPTX → * ============
            elif ext in ['pptx', 'ppt']:
                if target_format == 'pdf':
                    from tool_ppt_to_pdf import handler as pp_handler
                    result_str = pp_handler(ppt_path=source_path, output_path=output_path)
                    return json.loads(result_str)
                elif target_format == 'docx':
                    from tool_ppt_to_docx import handler as pd_handler
                    result_str = pd_handler(ppt_path=source_path, output_path=output_path)
                    return json.loads(result_str)
                elif target_format in ['png', 'jpg']:
                    from tool_ppt_to_images import handler as pi_handler
                    result_str = pi_handler(ppt_path=source_path, output_dir=output_path if os.path.isdir(output_path) else str(Path(output_path).parent / Path(output_path).stem), format=target_format)
                    return json.loads(result_str)
                elif target_format in ['md', 'txt']:
                    from pptx import Presentation
                    prs = Presentation(source_path)
                    lines = [f"# {Path(source_path).stem}\n"]
                    for i, slide in enumerate(prs.slides, 1):
                        lines.append(f"\n## Slide {i}\n")
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    if para.text.strip():
                                        lines.append(para.text + "\n")
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write('\n'.join(lines))
                else:
                    return {"success": False, "message": f"不支持 PPT→{target_format}"}

            # ============ 图片 → * ============
            elif ext in ['png', 'jpg', 'jpeg', 'bmp', 'tiff', 'webp']:
                if target_format == 'pdf':
                    try:
                        import img2pdf
                        with open(output_path, 'wb') as f:
                            f.write(img2pdf.convert(source_path))
                    except:
                        from PIL import Image
                        Image.open(source_path).convert('RGB').save(output_path, 'PDF')
                elif target_format in ['txt', 'md']:
                    import pytesseract
                    from PIL import Image
                    text = pytesseract.image_to_string(Image.open(source_path), lang='chi_sim+eng')
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(text)
                elif target_format == 'docx':
                    from docx import Document
                    from PIL import Image
                    doc = Document()
                    doc.add_picture(source_path)
                    # 同时尝试OCR
                    try:
                        import pytesseract
                        text = pytesseract.image_to_string(Image.open(source_path), lang='chi_sim+eng')
                        if text.strip():
                            doc.add_paragraph("OCR识别文本:")
                            doc.add_paragraph(text)
                    except:
                        pass
                    doc.save(output_path)
                elif target_format == 'xlsx':
                    import pytesseract
                    from PIL import Image
                    import pandas as pd
                    text = pytesseract.image_to_string(Image.open(source_path), lang='chi_sim+eng')
                    lines = [l for l in text.split('\n') if l.strip()]
                    pd.DataFrame(lines, columns=["OCR文本"]).to_excel(output_path, index=False, engine='openpyxl')
                elif target_format == 'jpg':
                    from PIL import Image
                    Image.open(source_path).convert('RGB').save(output_path, 'JPEG', quality=95)
                elif target_format == 'png':
                    from PIL import Image
                    Image.open(source_path).save(output_path, 'PNG')
                else:
                    return {"success": False, "message": f"不支持 图片→{target_format}"}

            # ============ Markdown → * ============
            elif ext == 'md':
                if target_format == 'pdf':
                    from tool_markdown_to_pdf import handler as mp_handler
                    result_str = mp_handler(md_path=source_path, output_path=output_path)
                    return json.loads(result_str)
                elif target_format == 'docx':
                    from tool_markdown_to_word import handler as mw_handler
                    result_str = mw_handler(md_path=source_path, output_path=output_path)
                    return json.loads(result_str)
                elif target_format == 'txt':
                    with open(source_path, 'r', encoding='utf-8') as f:
                        text = f.read()
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(text)
                elif target_format == 'html':
                    import markdown
                    with open(source_path, 'r', encoding='utf-8') as f:
                        md_text = f.read()
                    html = f"<html><body>\n{markdown.markdown(md_text, extensions=['extra','codehilite','tables'])}\n</body></html>"
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(html)
                else:
                    return {"success": False, "message": f"不支持 Markdown→{target_format}"}

            # ============ TXT → * ============
            elif ext == 'txt':
                if target_format == 'pdf':
                    from tool_txt_to_pdf import handler as tp_handler
                    result_str = tp_handler(txt_path=source_path, output_path=output_path)
                    return json.loads(result_str)
                elif target_format == 'docx':
                    from docx import Document
                    doc = Document()
                    with open(source_path, 'r', encoding='utf-8', errors='replace') as f:
                        for line in f:
                            if line.strip():
                                doc.add_paragraph(line.strip())
                    doc.save(output_path)
                elif target_format == 'md':
                    with open(source_path, 'r', encoding='utf-8', errors='replace') as f:
                        text = f.read()
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(text)
                elif target_format == 'xlsx':
                    import pandas as pd
                    with open(source_path, 'r', encoding='utf-8', errors='replace') as f:
                        lines = [l.strip() for l in f if l.strip()]
                    pd.DataFrame(lines, columns=["内容"]).to_excel(output_path, index=False, engine='openpyxl')
                else:
                    return {"success": False, "message": f"不支持 TXT→{target_format}"}

            # ============ HTML → * ============
            elif ext in ['html', 'htm']:
                from tool_html_converter import handler as hc_handler
                result_str = hc_handler(html_path=source_path, output_format=target_format, output_path=output_path)
                return json.loads(result_str)

            # ============ EPUB → * ============
            elif ext == 'epub':
                from tool_epub_converter import handler as ec_handler
                result_str = ec_handler(epub_path=source_path, output_format=target_format, output_path=output_path)
                return json.loads(result_str)

            # ============ CSV → * ============
            elif ext == 'csv':
                import pandas as pd
                df = pd.read_csv(source_path)
                if target_format == 'xlsx':
                    df.to_excel(output_path, index=False, engine='openpyxl')
                elif target_format == 'pdf':
                    from reportlab.lib.pagesizes import A4
                    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
                    from reportlab.lib import colors
                    doc = SimpleDocTemplate(output_path, pagesize=A4)
                    data = [df.columns.tolist()] + df.astype(str).values.tolist()
                    t = Table(data)
                    t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.grey),('GRID',(0,0),(-1,-1),0.5,colors.black)]))
                    doc.build([t])
                elif target_format == 'docx':
                    from docx import Document
                    doc = Document()
                    table = doc.add_table(rows=len(df)+1, cols=len(df.columns))
                    table.style = 'Light Grid Accent 1'
                    for j, col in enumerate(df.columns): table.rows[0].cells[j].text = str(col)
                    for i in range(len(df)):
                        for j in range(len(df.columns)): table.rows[i+1].cells[j].text = str(df.iloc[i,j])
                    doc.save(output_path)
                elif target_format == 'md':
                    lines = ["| " + " | ".join(str(c) for c in df.columns) + " |"]
                    lines.append("| " + " | ".join("---" for _ in df.columns) + " |")
                    for _, row in df.iterrows():
                        lines.append("| " + " | ".join(str(v) for v in row.values) + " |")
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write('\n'.join(lines))
                elif target_format == 'json':
                    df.to_json(output_path, orient='records', force_ascii=False, indent=2)
                elif target_format == 'txt':
                    df.to_csv(output_path, sep='\t', index=False)
                else:
                    return {"success": False, "message": f"不支持 CSV→{target_format}"}

            # ============ CAT格式 (TMX/TBX/XLIFF/SDLXLIFF) → * ============
            elif ext in ['tmx', 'tbx', 'xliff', 'xlf', 'sdlxliff']:
                from tool_cat_converter import handler as cc_handler
                result_str = cc_handler(cat_path=source_path, output_format=target_format, output_path=output_path)
                return json.loads(result_str)

            else:
                return {"success": False, "message": f"不支持的源格式: {ext}。支持: pdf/docx/xlsx/pptx/png/jpg/md/txt/csv/html/epub/tmx/tbx/xliff/sdlxliff"}

            # 最终检查
            if os.path.exists(output_path):
                size = os.path.getsize(output_path)
                return {"success": True, "output": output_path,
                        "message": f"✅ {ext.upper()}→{target_format.upper()} ({size/1024:.1f}KB): {output_path}"}
            else:
                return {"success": False, "output": "", "message": "输出文件未生成"}

        except Exception as e:
            import traceback
            return {"success": False, "output": "", "message": f"转换失败 [{ext}→{target_format}]: {str(e)}"}

    result = convert(**kwargs)
    return json.dumps(result, ensure_ascii=False)
